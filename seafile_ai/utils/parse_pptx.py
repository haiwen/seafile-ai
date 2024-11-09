import logging

from io import BytesIO
from operator import attrgetter

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE


logger = logging.getLogger(__name__)


def is_title(shape):
    if shape.is_placeholder and (
        shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
        or shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
        or shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE
        or shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE
    ):
        return True
    return False


def is_text_block(shape):
    if shape.has_text_frame:
        return True
    return False


def is_list_block(shape):
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def put_list(text, level):
    return '  ' * level + '* ' + text.strip() + '\n'


def put_table(table):
    table_text = ''
    gen_table_row = (
        lambda row: '| ' + ' | '.join([c.replace('\n', '<br />') for c in row]) + ' |'
    )
    table_text += gen_table_row(table[0]) + '\n'
    table_text += gen_table_row([':-:' for _ in table[0]]) + '\n'
    table_text += '\n'.join([gen_table_row(row) for row in table[1:]]) + '\n\n'
    return table_text


def process_title(shape, slide_idx):
    text = shape.text_frame.text.strip()
    return text


def process_text_block(shape):
    text = ''
    if is_list_block(shape):
        # generate list block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text += put_list(para.text, para.level)
    else:
        # generate paragraph block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text += para.text
    return text


def process_table(shape, _):
    table = [[cell.text for cell in row.cells] for row in shape.table.rows]
    if len(table) > 0:
        return put_table(table)
    else:
        return ''


def process_shapes(current_shapes, slide_id):
    slide_text = f'**SLIDE**: {slide_id+1} Content: \n'
    for idx, shape in enumerate(current_shapes):
        if is_title(shape):
            slide_text += shape.text_frame.text.strip() + '\n'
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            slide_text += process_table(shape, slide_id + 1) + '\n'
        elif shape.has_text_frame:
            slide_text += process_text_block(shape) + '\n'
    return slide_text + '-' * 25 + '\n'


def ungroup_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            res.extend(ungroup_shapes(shape.shapes))
        else:
            res.append(shape)
    return res


def get_pptx_text(file):
    prs = Presentation(BytesIO(file))
    pptx_text = ''
    for idx, slide in enumerate(prs.slides):
        try:
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
            pptx_text += process_shapes(shapes, idx)
        except Exception:
            continue
    return pptx_text
